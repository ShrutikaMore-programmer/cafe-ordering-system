"""
Patient Health Risk Assessment Report - PowerPoint Presentation Generator
Author: Shrutika More
Date: 2026-04-30
Course: Healthcare Data Analytics

This script creates a comprehensive presentation deck with the SQL query analysis results.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_DARK_BLUE = RGBColor(25, 55, 109)
COLOR_LIGHT_BLUE = RGBColor(65, 105, 225)
COLOR_ACCENT_GREEN = RGBColor(34, 177, 76)
COLOR_WARNING_RED = RGBColor(192, 0, 0)
COLOR_WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_DARK_BLUE
    title_shape.line.color.rgb = COLOR_DARK_BLUE
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(8.5), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, content in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK_BLUE
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_DARK_BLUE
    title_shape.line.color.rgb = COLOR_DARK_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4.5), Inches(0.4))
    left_header_frame = left_header.text_frame
    left_header_frame.text = left_title
    left_header_frame.paragraphs[0].font.size = Pt(16)
    left_header_frame.paragraphs[0].font.bold = True
    left_header_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT_BLUE
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, content in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_header = slide.shapes.add_textbox(Inches(5.25), Inches(1), Inches(4.5), Inches(0.4))
    right_header_frame = right_header.text_frame
    right_header_frame.text = right_title
    right_header_frame.paragraphs[0].font.size = Pt(16)
    right_header_frame.paragraphs[0].font.bold = True
    right_header_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT_BLUE
    
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, content in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# ==================== SLIDE 1: TITLE SLIDE ====================
add_title_slide(
    prs,
    "Patient Health Risk Assessment Report",
    "Healthcare Leadership Dashboard\n\nPresented by: Shrutika More\nDate: April 30, 2026"
)

# ==================== SLIDE 2: EXECUTIVE SUMMARY ====================
add_content_slide(
    prs,
    "Executive Summary",
    [
        "📊 Comprehensive Risk Assessment Framework",
        "✓ Analyzed patient population across all departments",
        "✓ Identified high-risk patients requiring immediate attention",
        "✓ Developed department-specific benchmarking metrics",
        "",
        "🎯 Key Business Impact:",
        "• Enables prioritized resource allocation",
        "• Supports proactive patient care management",
        "• Facilitates department-level performance comparison"
    ]
)

# ==================== SLIDE 3: REPORT OVERVIEW ====================
add_content_slide(
    prs,
    "Report Overview & Purpose",
    [
        "🏥 Business Question Answered:",
        "Which patients require priority care based on health risk factors?",
        "",
        "📈 Report Value:",
        "• Identifies 3 risk categories: High, Medium, Low",
        "• Provides ranking within risk category and department",
        "• Benchmarks individual performance against department averages",
        "",
        "🔍 Data Integration:",
        "Combines Patients, Medical Records, Lab Results, Doctors, and Departments"
    ]
)

# ==================== SLIDE 4: RISK CATEGORIZATION LOGIC ====================
add_two_column_slide(
    prs,
    "Risk Categorization Framework",
    "HIGH RISK Criteria",
    [
        "🔴 Patients Age 65+ with:",
        "  • 2+ chronic conditions",
        "  • 2+ abnormal lab results",
        "",
        "🔴 Any Patient with:",
        "  • 3+ chronic conditions",
        "",
        "⚠️  Priority: Immediate intervention required"
    ],
    "MEDIUM & LOW RISK",
    [
        "🟡 MEDIUM RISK Criteria:",
        "  • Age 50+ with 1+ conditions",
        "  • Exactly 2 chronic conditions",
        "  • 2+ abnormal lab results",
        "  • 5+ medical visits",
        "",
        "🟢 LOW RISK Criteria:",
        "  • All other patients with medical history",
        "  • Routine care management"
    ]
)

# ==================== SLIDE 5: WINDOW FUNCTIONS EXPLAINED ====================
add_content_slide(
    prs,
    "Advanced Analytics: Window Functions",
    [
        "🔧 RANK() OVER (PARTITION BY risk_category ORDER BY...)",
        "   → Identifies highest-risk patients within each risk group",
        "",
        "🔧 RANK() OVER (PARTITION BY department_name ORDER BY...)",
        "   → Enables department-specific prioritization",
        "",
        "🔧 AVG() OVER (PARTITION BY department_name)",
        "   → Calculates department benchmarks for comparison",
        "",
        "📊 Comparative Labels:",
        "   → Shows if patient is Above/Below/At department average"
    ]
)

# ==================== SLIDE 6: KEY FINDINGS - RISK DISTRIBUTION ====================
add_content_slide(
    prs,
    "Key Finding #1: Risk Distribution Analysis",
    [
        "📊 Patient Risk Category Breakdown:",
        "• HIGH RISK: Priority cases requiring immediate intervention",
        "• MEDIUM RISK: Patients requiring close monitoring",
        "• LOW RISK: Routine care patients",
        "",
        "💡 Key Insight:",
        "Risk stratification enables targeted resource allocation",
        "and proactive care management to prevent complications",
        "",
        "📌 Action Item:",
        "Develop specialized care protocols for high-risk patient cohorts"
    ]
)

# ==================== SLIDE 7: KEY FINDINGS - CHRONIC CONDITIONS ====================
add_content_slide(
    prs,
    "Key Finding #2: Chronic Condition Impact",
    [
        "🏥 Chronic Conditions Analysis:",
        "• Patients with 3+ conditions consistently in HIGH RISK",
        "• Age 65+ with 2+ conditions = elevated risk profile",
        "• Conditions like Hypertension, Diabetes, Heart Disease",
        "  are primary risk drivers",
        "",
        "💡 Key Insight:",
        "Multimorbidity (multiple chronic conditions) is strongest",
        "predictor of patient risk level",
        "",
        "📌 Action Item:",
        "Implement integrated care management for multimorbid patients"
    ]
)

# ==================== SLIDE 8: KEY FINDINGS - DEPARTMENT COMPARISON ====================
add_content_slide(
    prs,
    "Key Finding #3: Department Benchmarking",
    [
        "🏢 Department-Level Insights:",
        "• Significant variation in risk profiles across departments",
        "• Department averages enable fair performance comparison",
        "• Some departments handle higher-complexity patient loads",
        "",
        "💡 Key Insight:",
        "Department context essential for resource planning and",
        "staffing decisions - one-size-fits-all approach ineffective",
        "",
        "📌 Action Item:",
        "Tailor staffing and resource allocation by department risk profile"
    ]
)

# ==================== SLIDE 9: BUSINESS VALUE ====================
add_two_column_slide(
    prs,
    "Business Value & Strategic Impact",
    "Operational Benefits",
    [
        "✓ Prioritized Care Pathways",
        "   Focus resources on highest-risk patients first",
        "",
        "✓ Resource Optimization",
        "   Allocate staff based on department risk profiles",
        "",
        "✓ Performance Metrics",
        "   Department benchmarking identifies best practices",
        "",
        "✓ Risk Mitigation",
        "   Early intervention prevents adverse events"
    ],
    "Financial Benefits",
    [
        "💰 Cost Reduction",
        "   Prevent expensive ER visits and readmissions",
        "",
        "💰 Revenue Optimization",
        "   Improve case-mix management and billing accuracy",
        "",
        "💰 Outcome Improvement",
        "   Better patient outcomes = better quality metrics",
        "",
        "💰 Compliance",
        "   Demonstrate risk-based care delivery approach"
    ]
)

# ==================== SLIDE 10: RECOMMENDATIONS ====================
add_content_slide(
    prs,
    "Data-Driven Recommendations",
    [
        "1️⃣  Implement Risk-Based Care Protocols",
        "   Develop specialized care pathways for HIGH RISK patients",
        "",
        "2️⃣  Establish Department Care Coordinators",
        "   Assign specialists based on department risk profiles",
        "",
        "3️⃣  Create Automated Alerts & Monitoring",
        "   Alert clinical teams when patients reach risk thresholds",
        "",
        "4️⃣  Quarterly Review Cycle",
        "   Update risk assessments to track patient progression",
        "",
        "5️⃣  Integrate with EHR Systems",
        "   Real-time risk flags in patient records for all providers"
    ]
)

# ==================== SLIDE 11: IMPLEMENTATION ROADMAP ====================
add_content_slide(
    prs,
    "Implementation Roadmap",
    [
        "📅 Phase 1 (Month 1-2): Foundation",
        "   • Deploy SQL query in production database",
        "   • Train clinical staff on risk categories",
        "",
        "📅 Phase 2 (Month 3-4): Integration",
        "   • Connect to EHR systems",
        "   • Create automated alerts and dashboards",
        "",
        "📅 Phase 3 (Month 5-6): Optimization",
        "   • Gather feedback from clinical teams",
        "   • Refine risk thresholds based on outcomes",
        "",
        "📅 Phase 4 (Month 6+): Expansion",
        "   • Extend to predictive models",
        "   • Integrate social determinants of health"
    ]
)

# ==================== SLIDE 12: EXPECTED OUTCOMES ====================
add_two_column_slide(
    prs,
    "Expected Outcomes & Success Metrics",
    "Clinical Outcomes",
    [
        "📈 Reduce hospital readmissions",
        "   Target: 15-20% reduction in 90-day readmits",
        "",
        "📈 Improve patient safety",
        "   Target: Fewer adverse events in high-risk cohort",
        "",
        "📈 Better disease management",
        "   Target: Improved chronic condition control rates",
        "",
        "📈 Patient satisfaction",
        "   Target: Higher scores in proactive care perception"
    ],
    "Operational Metrics",
    [
        "🎯 Resource utilization",
        "   Measure staffing efficiency gains",
        "",
        "🎯 Time to intervention",
        "   Track time from risk identification to action",
        "",
        "🎯 Department compliance",
        "   Monitor adherence to risk protocols",
        "",
        "🎯 Data quality",
        "   Ensure accurate ongoing risk assessment"
    ]
)

# ==================== SLIDE 13: TECHNICAL SPECIFICATIONS ====================
add_content_slide(
    prs,
    "Technical Specifications & Data Integration",
    [
        "🗄️  Database Tables Integrated:",
        "   Patients → MedicalRecords → Doctors → Departments → LabResults",
        "",
        "📊 Output Columns (14 total):",
        "   Patient Demographics | Health Metrics | Risk Category",
        "   Risk Reason | Department Info | Rankings | Benchmarks",
        "",
        "⚙️  SQL Features Used:",
        "   CTEs (Common Table Expressions) | Window Functions",
        "   CASE Statements | Aggregation Functions | Joins",
        "",
        "🔄 Update Frequency: Monthly (adjustable based on needs)"
    ]
)

# ==================== SLIDE 14: CONCLUSION ====================
add_content_slide(
    prs,
    "Conclusion & Next Steps",
    [
        "✅ Comprehensive Risk Assessment Framework Complete",
        "",
        "🎯 Strategic Benefits:",
        "   • Prioritize high-risk patients for better outcomes",
        "   • Optimize resource allocation across departments",
        "   • Enable data-driven clinical decisions",
        "",
        "→ Next Steps:",
        "   1. Stakeholder review and approval (Leadership)",
        "   2. Implementation planning (IT & Clinical Teams)",
        "   3. Pilot deployment in 1-2 departments",
        "   4. Training and change management",
        "   5. Full rollout and monitoring",
        "",
        "📞 Questions & Discussion"
    ]
)

# Save presentation
output_path = "Patient_Health_Risk_Assessment_Report.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
